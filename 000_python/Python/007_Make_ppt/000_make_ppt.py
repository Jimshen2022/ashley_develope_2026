from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# 读取数据
df = excel_file.parse("Export")
df["Q'ty"] = pd.to_numeric(df["Q'ty"], errors="coerce").fillna(0)

# 按 Vendor 汇总超发数量
vendor_summary = df.groupby("Vendor")["Q'ty"].sum().sort_values(ascending=False)
top_vendors = vendor_summary.head(5)

# 新建 PPT
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]  # 标题布局

# 封面
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Vendor Over Shipment Report"
subtitle.text = "Vietnam Bonded Warehouse – Finished Furniture\nDate: May 2025"

# 数据总览
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Shipment Overview"

content = slide.placeholders[1]
total_records = len(df)
total_vendors = df["Vendor"].nunique()
total_qty = int(df["Q'ty"].sum())

content.text = (
    f"Total Over Shipment Records: {total_records}\n"
    f"Total Vendors Involved: {total_vendors}\n"
    f"Total Over Shipped Quantity: {total_qty}"
)

# Vendor 超发统计图表
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Top 5 Vendors by Over Shipped Quantity"

chart_data = CategoryChartData()
chart_data.categories = list(top_vendors.index)
chart_data.add_series("Over Shipped Qty", list(top_vendors.values))

x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
).chart

# 设置图表标题字体大小
chart.chart_title.has_text_frame = True
chart.chart_title.text_frame.text = "Top 5 Vendors with Most Over Shipments"

# 结束页
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Thank You"
slide.placeholders[1].text = "For inquiries, please contact the warehouse operations team."

# 保存
ppt_path = "/mnt/data/Vendor_Over_Shipment_Report.pptx"
prs.save(ppt_path)

ppt_path
